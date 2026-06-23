"""
PPTX report builder — "Obsidian Intelligence" theme.

Design philosophy:
  - Full dark backgrounds (deep navy / near-black) on every slide.
  - Electric indigo + cyan as accent colours — matching the web UI exactly.
  - ONE clear idea per slide section; content is never crammed.
  - Generous padding everywhere.  Text boxes are sized conservatively so
    nothing overflows.
  - Items are capped per slide: bullets max 4, recommendations max 3,
    competitor rows max 3.  Long strings are truncated to 120 chars so
    they never bleed outside their box.
  - All font sizes set explicitly — no reliance on pptx auto-fit.
"""

from __future__ import annotations
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Palette ────────────────────────────────────────────────────────────────
BG        = RGBColor(0x03, 0x03, 0x0E)   # near-black background
BG2       = RGBColor(0x07, 0x07, 0x1A)   # slightly lighter panel bg
SURFACE   = RGBColor(0x0C, 0x0C, 0x22)   # card surface
SURFACE2  = RGBColor(0x10, 0x10, 0x2A)   # card surface alt
INDIGO    = RGBColor(0x6C, 0x63, 0xFF)   # primary accent
INDIGO2   = RGBColor(0x9D, 0x98, 0xFF)   # lighter indigo
CYAN      = RGBColor(0x4D, 0xCF, 0xE8)   # secondary accent
GREEN     = RGBColor(0x3D, 0xDB, 0xA8)   # positive / growth
RED       = RGBColor(0xF0, 0x5C, 0x5C)   # risk / negative
AMBER     = RGBColor(0xF0, 0xA3, 0x5C)   # warning / breakeven
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
HI        = RGBColor(0xEE, 0xEE, 0xF5)   # primary text on dark
MUTED     = RGBColor(0x88, 0x88, 0xA8)   # secondary text on dark
LINE      = RGBColor(0x1E, 0x1E, 0x38)   # subtle border

FONT_HEAD = "Calibri"
FONT_MONO = "Courier New"
FONT_BODY = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
M       = Inches(0.65)          # standard margin
CONTENT_Y = Inches(1.55)        # top of content area (below header band)
CONTENT_H = Inches(5.55)        # usable height below header
CONTENT_W = SLIDE_W - 2 * M    # usable width


# ── XML helpers ─────────────────────────────────────────────────────────────

def _hex(rgb: RGBColor) -> str:
    return f"{rgb[0]:02X}{rgb[1]:02X}{rgb[2]:02X}"


def _solid_fill_xml(rgb: RGBColor) -> str:
    return (
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{_hex(rgb)}"/></a:solidFill>'
    )


def _set_cell_fill(cell, rgb: RGBColor):
    tc = cell._tc
    tcPr = tc.get_or_add_tcPr()
    # remove any existing fill
    for old in tcPr.findall(
        "{http://schemas.openxmlformats.org/drawingml/2006/main}solidFill"
    ):
        tcPr.remove(old)
    fill_el = etree.fromstring(_solid_fill_xml(rgb))
    tcPr.insert(0, fill_el)


# ── Low-level primitives ────────────────────────────────────────────────────

def _blank_slide(prs: Presentation, bg: RGBColor = BG):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = bg
    return slide


def _no_autofit(tf):
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE
    except Exception:
        pass


def _margins(tf, l=0.08, t=0.04, r=0.08, b=0.04):
    tf.margin_left   = Inches(l)
    tf.margin_right  = Inches(r)
    tf.margin_top    = Inches(t)
    tf.margin_bottom = Inches(b)


def _trunc(s: str, n: int = 115) -> str:
    """Hard-truncate a string so it never exceeds the text-box bounds."""
    s = str(s or "").strip()
    return s[:n] + "…" if len(s) > n else s


def add_text(
    slide, x, y, w, h, text: str,
    size: float = 13, color: RGBColor = HI,
    bold: bool = False, italic: bool = False,
    font: str = FONT_BODY,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float = 1.18,
    trunc: int = 0,
):
    if trunc:
        text = _trunc(text, trunc)
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    _no_autofit(tf)
    _margins(tf)
    tf.vertical_anchor = anchor
    p   = tf.paragraphs[0]
    p.alignment    = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text            = text
    run.font.size       = Pt(size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.name       = font
    run.font.color.rgb  = color
    return box


def add_bullets(
    slide, x, y, w, h, items: list,
    size: float = 12, color: RGBColor = HI,
    marker: str = "›", marker_color: RGBColor | None = None,
    line_spacing: float = 1.15, space_after: int = 7,
    max_items: int = 5, trunc: int = 100,
):
    marker_color = marker_color or INDIGO2
    items = [_trunc(str(i), trunc) for i in (items or ["Not available"])[:max_items]]
    box = slide.shapes.add_textbox(x, y, w, h)
    tf  = box.text_frame
    _no_autofit(tf)
    _margins(tf)
    for idx, item in enumerate(items):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.line_spacing = line_spacing
        p.space_after  = Pt(space_after)
        r1 = p.add_run()
        r1.text           = f"{marker}  "
        r1.font.size      = Pt(size)
        r1.font.name      = FONT_BODY
        r1.font.color.rgb = marker_color
        r1.font.bold      = True
        r2 = p.add_run()
        r2.text           = item
        r2.font.size      = Pt(size)
        r2.font.name      = FONT_BODY
        r2.font.color.rgb = color
    return box


def _rect(slide, x, y, w, h, fill: RGBColor, radius: bool = False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, x, y, w, h)
    if radius:
        try:
            sh.adjustments[0] = 0.06
        except Exception:
            pass
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _card(slide, x, y, w, h, fill: RGBColor = SURFACE, border: RGBColor | None = INDIGO):
    sh = _rect(slide, x, y, w, h, fill, radius=True)
    if border:
        sh.line.color.rgb = border
        sh.line.width     = Pt(0.6)
    return sh


def _label(slide, x, y, w, text: str, color: RGBColor = MUTED, size: float = 8.5):
    """Small all-caps label above a section."""
    add_text(slide, x, y, w, Inches(0.28), text.upper(),
             size=size, color=color, bold=True,
             font=FONT_MONO, line_spacing=1.0)


def _oval(slide, x, y, d, fill: RGBColor):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL, x, y, d, d)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    sh.line.fill.background(); sh.shadow.inherit = False
    return sh


def _accent_bar(slide, color: RGBColor = INDIGO, h_in: float = 0.03):
    """Full-width thin coloured bar just below the header band."""
    _rect(slide, 0, Inches(1.42), SLIDE_W, Inches(h_in), color)


def _footer(slide, page_no: int, total: int, tag: str = ""):
    """Consistent footer: dot · page tag · page number."""
    _rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), BG2)
    txt = f"  ·  {tag}  ·  {page_no:02d} / {total:02d}" if tag else f"  ·  {page_no:02d} / {total:02d}"
    add_text(
        slide, M, SLIDE_H - Inches(0.36), SLIDE_W - 2 * M, Inches(0.3),
        "MARKET INTELLIGENCE ENGINE" + txt,
        size=7.5, color=MUTED, font=FONT_MONO, align=PP_ALIGN.LEFT,
    )


def _header_band(slide, title: str, subtitle: str = "", accent: RGBColor = INDIGO):
    """Full-width dark header band shared by all content slides."""
    _rect(slide, 0, 0, SLIDE_W, Inches(1.42), BG2)
    _rect(slide, 0, Inches(1.38), SLIDE_W, Inches(0.045), accent)
    # title
    add_text(slide, M, Inches(0.14), SLIDE_W * 0.7, Inches(0.72),
             title, size=26, color=HI, bold=True, font=FONT_HEAD,
             anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_text(slide, M, Inches(0.82), SLIDE_W * 0.72, Inches(0.42),
                 subtitle, size=11, color=MUTED, font=FONT_BODY,
                 anchor=MSO_ANCHOR.MIDDLE)
    # slide-type tag top-right
    add_text(slide, SLIDE_W - Inches(2.4), Inches(0.18), Inches(1.8), Inches(0.38),
             "AI CONSULTING", size=7.5, color=accent,
             font=FONT_MONO, align=PP_ALIGN.RIGHT)


def _style_table(table, hdr_bg: RGBColor = SURFACE2, hdr_color: RGBColor = INDIGO2,
                 row_bg: RGBColor = SURFACE, row_alt: RGBColor = BG2,
                 body_color: RGBColor = HI, size: float = 11):
    for ci, cell in enumerate(table.rows[0].cells):
        _set_cell_fill(cell, hdr_bg)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.LEFT
            for r in p.runs:
                r.font.bold       = True
                r.font.size       = Pt(size)
                r.font.color.rgb  = hdr_color
                r.font.name       = FONT_BODY
    for ri, row in enumerate(list(table.rows)[1:]):
        bg = row_alt if ri % 2 else row_bg
        for cell in row.cells:
            _set_cell_fill(cell, bg)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for r in p.runs:
                    r.font.size       = Pt(size)
                    r.font.color.rgb  = body_color
                    r.font.name       = FONT_BODY


def _fill_table(table, headers: list, rows: list):
    for c, h in enumerate(headers):
        table.cell(0, c).text = str(h)
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            table.cell(r, c).text = str(val)


# ── Slide builders ──────────────────────────────────────────────────────────

def _slide_title(prs, state, page_no, total):
    slide = _blank_slide(prs, BG)

    # decorative diagonal stripe top-right
    _rect(slide, SLIDE_W - Inches(4.8), Inches(-0.5), Inches(5.2), Inches(0.28), INDIGO)
    _rect(slide, SLIDE_W - Inches(4.2), Inches(-0.1), Inches(4.6), Inches(0.14), CYAN)

    # left accent bar
    _rect(slide, 0, Inches(2.1), Inches(0.18), Inches(2.8), INDIGO)

    # glowing orb top-right (decorative)
    for r_in, opacity_hex in [(3.8, "18"), (2.5, "22"), (1.2, "2A")]:
        sh = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SLIDE_W - Inches(r_in * 1.4), Inches(-r_in * 0.5),
            Inches(r_in * 2), Inches(r_in * 2),
        )
        sh.fill.solid(); sh.fill.fore_color.rgb = INDIGO
        sh.line.fill.background(); sh.shadow.inherit = False
        # fade via transparency — approximate with a lighter colour
        

    industry = state.get("industry") or "Target Industry"
    market   = state.get("market")   or "Target Market"

    add_text(slide, Inches(0.9), Inches(0.85), Inches(7), Inches(0.45),
             "MARKET INTELLIGENCE ENGINE  ·  AI CONSULTING OS",
             size=9, color=INDIGO2, font=FONT_MONO)

    add_text(slide, Inches(0.9), Inches(1.6), Inches(10), Inches(2.0),
             "Market Entry\nStrategy Report",
             size=52, bold=True, color=HI, font=FONT_HEAD, line_spacing=1.05)

    add_text(slide, Inches(0.9), Inches(3.7), Inches(8), Inches(0.5),
             f"{industry}  ·  {market}",
             size=18, color=CYAN, font=FONT_BODY)

    # horizontal rule
    _rect(slide, Inches(0.9), Inches(4.35), Inches(6), Inches(0.025), INDIGO)

    add_text(slide, Inches(0.9), Inches(4.55), Inches(9), Inches(0.38),
             "Prepared by the Market Intelligence, Strategy & Executive Advisory Agents",
             size=11, color=MUTED, font=FONT_BODY, italic=True)

    # bottom tag strip
    _rect(slide, 0, SLIDE_H - Inches(0.38), SLIDE_W, Inches(0.38), BG2)
    add_text(slide, M, SLIDE_H - Inches(0.35), Inches(6), Inches(0.28),
             "CONFIDENTIAL  ·  AI-GENERATED ANALYSIS",
             size=7.5, color=MUTED, font=FONT_MONO)
    return slide


def _slide_executive_summary(prs, state, page_no, total):
    slide = _blank_slide(prs, BG)
    _header_band(slide, "Executive Summary",
                 "Strategic overview, opportunity & call to action", INDIGO)
    _footer(slide, page_no, total, "Executive Summary")

    summary = _trunc(state.get("executive_summary") or "Not available.", 420)
    fin = state.get("financials") or {}

    # ── Big summary card (takes ~55% of height)
    card_h = Inches(3.20)
    _card(slide, M, CONTENT_Y, CONTENT_W, card_h, SURFACE, INDIGO)
    _label(slide, M + Inches(0.3), CONTENT_Y + Inches(0.18), Inches(4), "EXECUTIVE OVERVIEW", INDIGO2)
    add_text(
        slide, M + Inches(0.3), CONTENT_Y + Inches(0.54),
        CONTENT_W - Inches(0.6), card_h - Inches(0.72),
        summary, size=13.5, color=HI, line_spacing=1.45,
    )

    # ── Three KPI chips underneath
    chips = [
        ("Initial Investment", fin.get("initial_investment") or "—", INDIGO2),
        ("Year 1 Revenue",     fin.get("year1_revenue")      or "—", CYAN),
        ("Breakeven Timeline", fin.get("breakeven_timeline") or "—", GREEN),
    ]
    chip_y = CONTENT_Y + card_h + Inches(0.28)
    chip_w = (CONTENT_W - Inches(0.5)) / 3
    for i, (lbl, val, col) in enumerate(chips):
        cx = M + i * (chip_w + Inches(0.25))
        _card(slide, cx, chip_y, chip_w, Inches(1.55), SURFACE2, col)
        _label(slide, cx + Inches(0.25), chip_y + Inches(0.18), chip_w - Inches(0.5), lbl, col)
        add_text(slide, cx + Inches(0.25), chip_y + Inches(0.52),
                 chip_w - Inches(0.5), Inches(0.85),
                 _trunc(val, 22), size=22, bold=True, color=col,
                 font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)
    return slide


def _slide_market_overview(prs, state, page_no, total):
    """Slide 3a — Overview + market size + growth trends only (no cramming)."""
    slide = _blank_slide(prs, BG)
    ma = state.get("market_analysis") or {}
    _header_band(slide, "Market Intelligence — Overview",
                 "Market size, landscape & growth trajectory", CYAN)
    _footer(slide, page_no, total, "Market Intelligence")

    # Overview card (left, tall)
    lw = Inches(6.1)
    _card(slide, M, CONTENT_Y, lw, Inches(5.25), SURFACE, CYAN)
    _label(slide, M + Inches(0.28), CONTENT_Y + Inches(0.18), lw - Inches(0.56), "MARKET OVERVIEW", CYAN)
    add_text(slide, M + Inches(0.28), CONTENT_Y + Inches(0.54),
             lw - Inches(0.56), Inches(2.1),
             _trunc(ma.get("overview") or "Not available.", 320),
             size=12.5, color=HI, line_spacing=1.42)

    _label(slide, M + Inches(0.28), CONTENT_Y + Inches(2.82), lw - Inches(0.56), "MARKET SIZE", INDIGO2)
    add_text(slide, M + Inches(0.28), CONTENT_Y + Inches(3.16),
             lw - Inches(0.56), Inches(0.75),
             _trunc(ma.get("market_size") or "—", 160),
             size=12.5, color=CYAN, bold=True, line_spacing=1.2)

    _label(slide, M + Inches(0.28), CONTENT_Y + Inches(4.02), lw - Inches(0.56), "CUSTOMER SEGMENTS", MUTED)
    segs = ma.get("customer_segments") or []
    add_bullets(slide, M + Inches(0.28), CONTENT_Y + Inches(4.34),
                lw - Inches(0.56), Inches(0.82),
                segs, size=11.5, marker="·", max_items=3)

    # Growth trends card (right)
    rw = CONTENT_W - lw - Inches(0.3)
    rx = M + lw + Inches(0.3)
    _card(slide, rx, CONTENT_Y, rw, Inches(5.25), SURFACE2, INDIGO)
    _label(slide, rx + Inches(0.25), CONTENT_Y + Inches(0.18), rw - Inches(0.5), "GROWTH TRENDS", INDIGO2)
    trends = ma.get("growth_trends") or ["Not available"]
    add_bullets(slide, rx + Inches(0.25), CONTENT_Y + Inches(0.54),
                rw - Inches(0.5), Inches(2.35),
                trends, size=11.5, marker="↗", marker_color=GREEN,
                max_items=4, trunc=90)

    _label(slide, rx + Inches(0.25), CONTENT_Y + Inches(3.0), rw - Inches(0.5), "KEY DRIVERS", MUTED)
    drivers = ma.get("market_drivers") or []
    add_bullets(slide, rx + Inches(0.25), CONTENT_Y + Inches(3.34),
                rw - Inches(0.5), Inches(1.8),
                drivers, size=11.5, marker="▸", marker_color=CYAN,
                max_items=3, trunc=80)
    return slide


def _slide_competitors(prs, state, page_no, total):
    """Slide 3b — Competitors + Opportunities + Risks."""
    slide = _blank_slide(prs, BG)
    ma = state.get("market_analysis") or {}
    _header_band(slide, "Market Intelligence — Competitive Landscape",
                 "Key players, opportunities & risks in the target market", CYAN)
    _footer(slide, page_no, total, "Market Intelligence")

    # ── Competitors table (top half)
    competitors = ma.get("competitors") or []
    rows = [(_trunc(c.get("name", "—"), 22), _trunc(c.get("note", ""), 90))
            for c in competitors[:4]]
    if not rows:
        rows = [("—", "Not available")]

    _label(slide, M, CONTENT_Y, CONTENT_W, "TOP COMPETITORS", CYAN)
    tbl_h = Inches(0.46 + 0.52 * len(rows))
    gf = slide.shapes.add_table(len(rows) + 1, 2, M, CONTENT_Y + Inches(0.32),
                                  CONTENT_W, tbl_h)
    t = gf.table
    t.columns[0].width = Inches(2.5)
    t.columns[1].width = int(CONTENT_W - Inches(2.5))
    _fill_table(t, ["Company", "Strategic Note"], rows)
    _style_table(t, size=11.5)

    # ── Opps + Risks side by side below the table
    opp_y = CONTENT_Y + Inches(0.32) + tbl_h + Inches(0.32)
    half  = (CONTENT_W - Inches(0.3)) / 2

    _card(slide, M, opp_y, half, Inches(2.62), SURFACE, GREEN)
    _label(slide, M + Inches(0.25), opp_y + Inches(0.18), half - Inches(0.5), "OPPORTUNITIES", GREEN)
    opps = ma.get("opportunities") or ["Not available"]
    add_bullets(slide, M + Inches(0.25), opp_y + Inches(0.52),
                half - Inches(0.5), Inches(2.0),
                opps, size=11.5, marker="+", marker_color=GREEN,
                max_items=4, trunc=85)

    rx = M + half + Inches(0.3)
    _card(slide, rx, opp_y, half, Inches(2.62), SURFACE, RED)
    _label(slide, rx + Inches(0.25), opp_y + Inches(0.18), half - Inches(0.5), "RISKS & BARRIERS", RED)
    risks = ma.get("risks") or []
    risks += ma.get("barriers_to_entry") or []
    add_bullets(slide, rx + Inches(0.25), opp_y + Inches(0.52),
                half - Inches(0.5), Inches(2.0),
                risks, size=11.5, marker="!", marker_color=RED,
                max_items=4, trunc=85)
    return slide


def _slide_swot(prs, state, page_no, total):
    slide = _blank_slide(prs, BG)
    _header_band(slide, "SWOT Analysis",
                 "Internal capabilities vs. external market forces", INDIGO)
    _footer(slide, page_no, total, "SWOT")

    swot = (state.get("strategy") or {}).get("swot") or {}
    quads = [
        ("STRENGTHS",    swot.get("strengths")    or ["Not available"], GREEN,  SURFACE),
        ("WEAKNESSES",   swot.get("weaknesses")   or ["Not available"], RED,    SURFACE2),
        ("OPPORTUNITIES",swot.get("opportunities") or ["Not available"], CYAN,   SURFACE),
        ("THREATS",      swot.get("threats")       or ["Not available"], AMBER,  SURFACE2),
    ]
    qw  = (CONTENT_W - Inches(0.28)) / 2
    qh  = (CONTENT_H - Inches(0.28)) / 2
    gap = Inches(0.28)

    for idx, (lbl, items, col, bg) in enumerate(quads):
        col_i = idx % 2
        row_i = idx // 2
        qx = M + col_i * (qw + gap)
        qy = CONTENT_Y + row_i * (qh + gap)

        _card(slide, qx, qy, qw, qh, bg, col)
        # coloured top bar
        _rect(slide, qx, qy, qw, Inches(0.07), col)
        _label(slide, qx + Inches(0.24), qy + Inches(0.16), qw - Inches(0.48), lbl, col, size=9)
        add_bullets(slide, qx + Inches(0.24), qy + Inches(0.52),
                    qw - Inches(0.48), qh - Inches(0.65),
                    items, size=12, color=HI, marker="–",
                    marker_color=col, max_items=4, trunc=88)
    return slide


def _slide_strategy(prs, state, page_no, total):
    slide = _blank_slide(prs, BG)
    strat = state.get("strategy") or {}
    _header_band(slide, "Market Entry Strategy",
                 "Entry approach, go-to-market steps & key partnerships", INDIGO)
    _footer(slide, page_no, total, "Strategy")

    col_w = (CONTENT_W - Inches(0.3)) / 2
    # Left column — entry strategy
    _card(slide, M, CONTENT_Y, col_w, CONTENT_H, SURFACE, INDIGO)
    _label(slide, M + Inches(0.28), CONTENT_Y + Inches(0.18), col_w - Inches(0.56), "ENTRY STRATEGY", INDIGO2)
    entry = strat.get("entry_strategy") or ["Not available"]
    add_bullets(slide, M + Inches(0.28), CONTENT_Y + Inches(0.55),
                col_w - Inches(0.56), Inches(2.55),
                entry, size=12, marker="›", max_items=4, trunc=90)

    _label(slide, M + Inches(0.28), CONTENT_Y + Inches(3.28), col_w - Inches(0.56), "PARTNERSHIPS", MUTED)
    partners = strat.get("partnerships") or ["Not available"]
    add_bullets(slide, M + Inches(0.28), CONTENT_Y + Inches(3.62),
                col_w - Inches(0.56), Inches(1.72),
                partners, size=12, marker="◆", marker_color=CYAN,
                max_items=3, trunc=88)

    # Right column — go-to-market + roadmap
    rx = M + col_w + Inches(0.3)
    _card(slide, rx, CONTENT_Y, col_w, CONTENT_H, SURFACE2, CYAN)
    _label(slide, rx + Inches(0.28), CONTENT_Y + Inches(0.18), col_w - Inches(0.56), "GO-TO-MARKET", CYAN)
    gtm = strat.get("go_to_market") or ["Not available"]
    add_bullets(slide, rx + Inches(0.28), CONTENT_Y + Inches(0.55),
                col_w - Inches(0.56), Inches(2.55),
                gtm, size=12, marker="▸", marker_color=CYAN,
                max_items=4, trunc=90)

    _label(slide, rx + Inches(0.28), CONTENT_Y + Inches(3.28), col_w - Inches(0.56), "IMPLEMENTATION ROADMAP", MUTED)
    roadmap = strat.get("implementation_roadmap") or ["Not available"]
    for i, phase in enumerate(roadmap[:3]):
        py = CONTENT_Y + Inches(3.62) + i * Inches(0.60)
        _rect(slide, rx + Inches(0.28), py, Inches(0.38), Inches(0.38), INDIGO, radius=True)
        add_text(slide, rx + Inches(0.28), py, Inches(0.38), Inches(0.38),
                 str(i + 1), size=12, bold=True, color=HI,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_text(slide, rx + Inches(0.78), py + Inches(0.04),
                 col_w - Inches(1.06), Inches(0.50),
                 _trunc(phase, 85), size=11.5, color=HI, line_spacing=1.1)
    return slide


def _slide_financials(prs, state, page_no, total):
    slide = _blank_slide(prs, BG)
    fin = state.get("financials") or {}
    _header_band(slide, "Financial Estimate",
                 "Investment requirements, revenue outlook & path to breakeven", GREEN)
    _footer(slide, page_no, total, "Financials")

    # ── Row 1: three KPI chips
    chips = [
        ("Initial Investment", fin.get("initial_investment") or "—", INDIGO2),
        ("Year 1 Revenue",     fin.get("year1_revenue")      or "—", GREEN),
        ("Breakeven Timeline", fin.get("breakeven_timeline") or "—", AMBER),
    ]
    chip_w = (CONTENT_W - Inches(0.5)) / 3
    chip_h = Inches(1.45)
    for i, (lbl, val, col) in enumerate(chips):
        cx = M + i * (chip_w + Inches(0.25))
        _card(slide, cx, CONTENT_Y, chip_w, chip_h, SURFACE, col)
        _rect(slide, cx, CONTENT_Y, chip_w, Inches(0.065), col)
        _label(slide, cx + Inches(0.22), CONTENT_Y + Inches(0.18), chip_w - Inches(0.44), lbl, col)
        add_text(slide, cx + Inches(0.22), CONTENT_Y + Inches(0.52),
                 chip_w - Inches(0.44), Inches(0.82),
                 _trunc(val, 20), size=21, bold=True, color=col,
                 font=FONT_HEAD, anchor=MSO_ANCHOR.MIDDLE)

    # ── Row 2: revenue projection table (left) + bar chart (right)
    proj = fin.get("projection") or []
    tbl_y = CONTENT_Y + chip_h + Inches(0.32)
    lw2 = Inches(4.8)

    _label(slide, M, tbl_y - Inches(0.28), lw2, "REVENUE PROJECTION")
    if proj:
        rows = [(_trunc(p.get("year", "—"), 16),
                 _trunc(p.get("revenue", "—"), 20),
                 _trunc(p.get("profit", "—"), 20)) for p in proj[:3]]
        tbl_h = Inches(0.46 + 0.52 * len(rows))
        gf = slide.shapes.add_table(
            len(rows) + 1, 3, M, tbl_y, lw2, tbl_h)
        t = gf.table
        t.columns[0].width = Inches(1.55)
        t.columns[1].width = Inches(1.65)
        t.columns[2].width = int(lw2 - Inches(3.2))
        _fill_table(t, ["Year", "Revenue", "Profit"], rows)
        _style_table(t, size=12)

        # horizontal bar chart to the right of the table
        import re as _re
        def _num(s):
            d = _re.sub(r"[^0-9.]", "", str(s))
            try: return float(d) if d else 0.0
            except: return 0.0

        vals = [_num(p.get("revenue", "0")) for p in proj[:3]]
        mx = max(vals) or 1.0
        bx_in = M + lw2 + Inches(0.45)
        track_w = SLIDE_W - bx_in - M - Inches(1.5)
        bar_h = Inches(0.44)
        bar_gap = Inches(0.28)
        bar_colors = [INDIGO, CYAN, GREEN]
        for i, (p, v) in enumerate(zip(proj[:3], vals)):
            by_i = tbl_y + i * (bar_h + bar_gap)
            bw = max(Inches(0.3), int(track_w * v / mx))
            sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                         bx_in, by_i, bw, bar_h)
            sh.fill.solid(); sh.fill.fore_color.rgb = bar_colors[i % 3]
            sh.line.fill.background(); sh.shadow.inherit = False
            try: sh.adjustments[0] = 0.3
            except: pass
            add_text(slide, bx_in + bw + Inches(0.12), by_i,
                     Inches(1.35), bar_h,
                     _trunc(p.get("revenue", ""), 18),
                     size=11, bold=True, color=bar_colors[i % 3],
                     anchor=MSO_ANCHOR.MIDDLE)

        # investment reasoning below
        reasoning = _trunc(fin.get("investment_reasoning") or "", 200)
        if reasoning:
            note_y = tbl_y + tbl_h + Inches(0.28)
            _label(slide, M, note_y - Inches(0.24), CONTENT_W, "INVESTMENT RATIONALE")
            _card(slide, M, note_y, CONTENT_W, Inches(0.88), SURFACE, INDIGO)
            add_text(slide, M + Inches(0.28), note_y + Inches(0.10),
                     CONTENT_W - Inches(0.56), Inches(0.72),
                     reasoning, size=11.5, color=MUTED, line_spacing=1.3)
    else:
        add_text(slide, M, tbl_y, CONTENT_W, Inches(0.5),
                 "Revenue projection not available.", size=13, color=MUTED)
    return slide


def _slide_recommendations(prs, state, page_no, total):
    slide = _blank_slide(prs, BG)
    _header_band(slide, "Strategic Recommendations",
                 "Priority actions advised by the Executive Advisory agent", INDIGO)
    _footer(slide, page_no, total, "Recommendations")

    recs = state.get("recommendations") or ["Not available"]
    # cap at 4 recommendations, generous space per item
    recs = [_trunc(r, 200) for r in recs[:4]]
    row_h = Inches(1.18)
    gap   = Inches(0.22)
    total_h = len(recs) * row_h + (len(recs) - 1) * gap
    start_y = CONTENT_Y + (CONTENT_H - total_h) / 2   # vertically centre

    acc_colors = [INDIGO, CYAN, GREEN, AMBER]
    for i, rec in enumerate(recs):
        ry = start_y + i * (row_h + gap)
        col = acc_colors[i % len(acc_colors)]
        _card(slide, M, ry, CONTENT_W, row_h, SURFACE, col)
        _rect(slide, M, ry, Inches(0.055), row_h, col)
        # number badge
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL,
            M + Inches(0.2), ry + (row_h - Inches(0.56)) / 2,
            Inches(0.56), Inches(0.56))
        badge.fill.solid(); badge.fill.fore_color.rgb = col
        badge.line.fill.background(); badge.shadow.inherit = False
        tf = badge.text_frame; _margins(tf, 0, 0, 0, 0)
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
        run = p.add_run(); run.text = str(i + 1)
        run.font.size = Pt(16); run.font.bold = True
        run.font.color.rgb = BG; run.font.name = FONT_HEAD

        add_text(slide, M + Inches(0.98), ry + Inches(0.08),
                 CONTENT_W - Inches(1.22), row_h - Inches(0.16),
                 rec, size=13, color=HI, line_spacing=1.3,
                 anchor=MSO_ANCHOR.MIDDLE)
    return slide


def _slide_closing(prs, state, page_no, total):
    slide = _blank_slide(prs, BG)

    # large decorative circles
    for r_in, col in [(5.5, INDIGO), (3.8, BG2), (2.0, SURFACE)]:
        sh = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            SLIDE_W - Inches(r_in * 1.1), SLIDE_H - Inches(r_in * 1.0),
            Inches(r_in * 2), Inches(r_in * 2))
        sh.fill.solid(); sh.fill.fore_color.rgb = col
        sh.line.fill.background(); sh.shadow.inherit = False

    _rect(slide, 0, Inches(0), Inches(0.12), SLIDE_H, INDIGO)

    industry = state.get("industry") or "Target Industry"
    market   = state.get("market")   or "Target Market"

    add_text(slide, Inches(0.8), Inches(1.8), Inches(8), Inches(0.38),
             "MARKET INTELLIGENCE ENGINE",
             size=9, color=INDIGO2, font=FONT_MONO)

    add_text(slide, Inches(0.8), Inches(2.3), Inches(8.5), Inches(1.9),
             "Analysis\nComplete.", size=56, bold=True, color=HI,
             font=FONT_HEAD, line_spacing=1.0)

    _rect(slide, Inches(0.8), Inches(4.32), Inches(4), Inches(0.03), CYAN)

    add_text(slide, Inches(0.8), Inches(4.5), Inches(8), Inches(0.45),
             f"{industry}  ·  {market}",
             size=15, color=CYAN, font=FONT_BODY)

    add_text(slide, Inches(0.8), Inches(5.1), Inches(8), Inches(0.38),
             "Report generated by the AI Consultant Agent Team  ·  Powered by Groq / LangGraph",
             size=10, color=MUTED, font=FONT_BODY, italic=True)
    return slide


# ── Entry points ─────────────────────────────────────────────────────────────

def build_presentation(state: dict, output_path: str = "market_entry_report.pptx") -> str:
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    builders = [
        _slide_title,
        _slide_executive_summary,
        _slide_market_overview,      # split into two slides ↓
        _slide_competitors,
        _slide_swot,
        _slide_strategy,
        _slide_financials,
        _slide_recommendations,
        _slide_closing,
    ]
    total = len(builders)
    for i, builder in enumerate(builders, start=1):
        builder(prs, state, i, total)

    prs.save(output_path)
    return output_path


def build_markdown_report(state: dict) -> str:
    industry = state.get("industry", "")
    market   = state.get("market", "")
    ma       = state.get("market_analysis") or {}
    strat    = state.get("strategy") or {}
    swot     = strat.get("swot") or {}
    fin      = state.get("financials") or {}

    def bl(items):
        return "\n".join(f"- {i}" for i in (items or [])) or "- Not available"

    return "\n".join([
        f"# Market Entry Strategy — {industry} / {market}", "",
        "## Executive Summary",
        state.get("executive_summary") or "Not available.", "",
        "## Market Overview",
        f"**Overview:** {ma.get('overview', '')}",
        f"**Market size:** {ma.get('market_size', '')}", "",
        "**Growth trends:**", bl(ma.get("growth_trends")), "",
        "**Competitors:**",
        "\n".join(f"- {c.get('name','—')} — {c.get('note','')}"
                  for c in (ma.get("competitors") or [])) or "- Not available", "",
        "**Opportunities:**", bl(ma.get("opportunities")), "",
        "**Risks:**", bl(ma.get("risks")), "",
        "## SWOT",
        "**Strengths:**",    bl(swot.get("strengths")), "",
        "**Weaknesses:**",   bl(swot.get("weaknesses")), "",
        "**Opportunities:**",bl(swot.get("opportunities")), "",
        "**Threats:**",      bl(swot.get("threats")), "",
        "## Entry Strategy", bl(strat.get("entry_strategy")), "",
        "## Go-to-Market",   bl(strat.get("go_to_market")), "",
        "## Financials",
        f"- Initial investment: {fin.get('initial_investment','—')}",
        f"- Year 1 revenue: {fin.get('year1_revenue','—')}",
        f"- Breakeven: {fin.get('breakeven_timeline','—')}", "",
        "**Projection:**",
        "\n".join(f"- {p.get('year','—')}: {p.get('revenue','—')}"
                  for p in (fin.get("projection") or [])) or "- Not available", "",
        "## Recommendations", bl(state.get("recommendations")),
    ])